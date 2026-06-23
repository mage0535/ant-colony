from __future__ import annotations

import os
import subprocess


class OfficeCliProvider:
    """Local OfficeCLI provider for private Office document processing."""

    def __init__(self) -> None:
        from src.tools.document_tool import OFFICECLI

        self.binary = OFFICECLI

    def is_configured(self) -> bool:
        return os.path.isfile(self.binary)

    def healthcheck(self) -> str | None:
        if not self.is_configured():
            return None
        try:
            result = subprocess.run(
                [self.binary, "--version"],
                capture_output=True,
                text=True,
                timeout=15,
                check=False,
            )
            if result.returncode != 0:
                return result.stderr.strip() or "officecli unavailable"
            return result.stdout.strip() or "officecli-ready"
        except Exception as exc:
            return f"OfficeCLI healthcheck failed: {exc}"

    def healthcheck_office(self) -> str | None:
        return self.healthcheck()

    def generate_docx_document(self, title: str, content: str, template_path: str | None = None) -> str | None:
        from src.tools.document_tool import generate_report

        return generate_report(title, content, "docx", template_path=template_path)

    def generate_xlsx_document(self, title: str, content: str, template_path: str | None = None) -> str | None:
        from src.tools.document_tool import generate_report

        return generate_report(title, content, "xlsx", template_path=template_path)

    def generate_pptx_document(self, title: str, content: str, template_path: str | None = None) -> str | None:
        from src.tools.document_tool import generate_report

        return generate_report(title, content, "pptx", template_path=template_path)

    def extract_docx_template_outline(self, path: str):
        from src.tools.document_tool import extract_docx_template_outline

        return extract_docx_template_outline(path)

    def extract_xlsx_template_outline(self, path: str):
        from src.tools.document_tool import extract_xlsx_template_outline

        return extract_xlsx_template_outline(path)

    def extract_pptx_template_outline(self, path: str):
        from src.tools.document_tool import extract_pptx_template_outline

        return extract_pptx_template_outline(path)

    def read_docx_document(self, path: str) -> str | None:
        from docx import Document

        doc = Document(path)
        lines = [p.text.strip() for p in doc.paragraphs if (p.text or "").strip()]
        return "\n".join(lines) if lines else ""

    def read_xlsx_document(self, path: str) -> str | None:
        from openpyxl import load_workbook

        wb = load_workbook(path, data_only=True)
        lines: list[str] = []
        for ws in wb.worksheets:
            lines.append(f"[Sheet] {ws.title}")
            for row in ws.iter_rows(values_only=True):
                values = [str(v) for v in row if v is not None and str(v).strip()]
                if values:
                    lines.append("\t".join(values))
        return "\n".join(lines) if lines else ""

    def read_pptx_document(self, path: str) -> str | None:
        from pptx import Presentation

        prs = Presentation(path)
        lines: list[str] = []
        for idx, slide in enumerate(prs.slides, start=1):
            lines.append(f"[Slide {idx}]")
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False):
                    text = getattr(shape, "text", "").strip()
                    if text:
                        lines.append(text)
        return "\n".join(lines) if lines else ""
