"""Document generation via OfficeCLI — DOCX reports, XLSX exports.
Uses the correct OfficeCLI API:
  - add (not set --prop text=) for adding new paragraphs/content
  - set for modifying existing properties (like cell values in xlsx)
"""
from __future__ import annotations

import logging
import os
import re
import shutil
import subprocess
import uuid
from copy import deepcopy
from typing import Any

logger = logging.getLogger(__name__)

OFFICECLI = "/usr/local/bin/officecli"
DOCUMENTS_DIR = os.path.join(os.path.dirname(os.path.dirname(os.path.dirname(__file__))), "data", "documents")


def _cli(args: list[str]) -> tuple[int, str]:
    cmd = [OFFICECLI] + args
    result = subprocess.run(cmd, capture_output=True, text=True, timeout=30)
    return result.returncode, (result.stderr or result.stdout)


def generate_report(title: str, content: str, format: str = "docx", template_path: str | None = None) -> str:
    """Generate a document using OfficeCLI."""
    title = title.replace("\x00", "")
    content = content.replace("\x00", "")
    logger.info("generate_report: title=%s content_len=%d format=%s", title, len(content), format)
    allowed = {"docx", "xlsx", "pptx"}
    if format not in allowed:
        return f"不支持的格式: {format}，可选: {', '.join(allowed)}"
    if not os.path.isfile(OFFICECLI):
        return "OfficeCLI 未安装"
    if len(content.strip()) < 5:
        return "文档内容太短，请提供更详细的文档内容后再生成。"

    os.makedirs(DOCUMENTS_DIR, exist_ok=True)
    filename = f"{_sanitize_filename(title)}.{format}"
    filepath = os.path.join(DOCUMENTS_DIR, filename)
    if os.path.exists(filepath):
        stem, ext = os.path.splitext(filename)
        filepath = os.path.join(DOCUMENTS_DIR, f"{stem}_{uuid.uuid4().hex[:8]}{ext}")

    try:
        if format == "docx":
            if template_path and os.path.isfile(template_path):
                _build_docx_from_template(template_path, filepath, content)
            else:
                rc, out = _cli(["create", filepath])
                if rc != 0:
                    return f"创建文档失败: {out}"
                _build_docx(filepath, content)
        elif format == "xlsx":
            if template_path and os.path.isfile(template_path):
                _build_xlsx_from_template(template_path, filepath, content)
            else:
                rc, out = _cli(["create", filepath])
                if rc != 0:
                    return f"创建文档失败: {out}"
                _build_xlsx(filepath, content)
        elif format == "pptx":
            if template_path and os.path.isfile(template_path):
                _build_pptx_from_template(template_path, filepath, content)
            else:
                rc, out = _cli(["create", filepath])
                if rc != 0:
                    return f"创建文档失败: {out}"
                _build_pptx(filepath, content)

        size = os.path.getsize(filepath)
        logger.info("Document created: %s (%d bytes)", filepath, size)
        return filepath

    except Exception as e:
        logger.exception("Document generation failed")
        return f"文档生成失败: {e}"


def extract_docx_template_outline(filepath: str) -> dict[str, Any]:
    """Extract a stable outline from a DOCX template for later structured filling."""
    try:
        from docx import Document
    except ImportError:
        logger.warning("python-docx not installed; outline extraction skipped")
        return {"paragraphs": [], "tables": []}

    try:
        doc = Document(filepath)
    except Exception as exc:
        logger.warning("Failed to open DOCX template %s: %s", filepath, exc)
        return {"paragraphs": [], "tables": []}

    paragraphs: list[dict[str, Any]] = []
    for index, para in enumerate(doc.paragraphs):
        text = (para.text or "").strip()
        if not text:
            continue
        style_name = getattr(getattr(para, "style", None), "name", "") or "Normal"
        anchor = text if style_name.startswith("Heading") else text[:40]
        paragraphs.append(
            {
                "index": index,
                "style": style_name,
                "text": text,
                "anchor": anchor,
            }
        )

    tables: list[dict[str, Any]] = []
    for table_index, table in enumerate(doc.tables):
        row_values: list[list[str]] = []
        max_cols = 0
        for row in table.rows:
            cell_values = [(cell.text or "").strip() for cell in row.cells]
            row_values.append(cell_values)
            max_cols = max(max_cols, len(cell_values))
        tables.append(
            {
                "index": table_index,
                "rows": len(row_values),
                "cols": max_cols,
                "cells": row_values,
            }
        )

    return {"paragraphs": paragraphs, "tables": tables}


def extract_xlsx_template_outline(filepath: str) -> dict[str, Any]:
    try:
        from openpyxl import load_workbook
    except ImportError:
        logger.warning("openpyxl not installed; xlsx outline extraction skipped")
        return {"sheets": []}

    try:
        wb = load_workbook(filepath, data_only=True)
    except Exception as exc:
        logger.warning("Failed to open XLSX template %s: %s", filepath, exc)
        return {"sheets": []}

    sheets: list[dict[str, Any]] = []
    for ws in wb.worksheets:
        rows: list[list[str]] = []
        max_cols = 0
        for row in ws.iter_rows(values_only=True):
            values = ["" if value is None else str(value).strip() for value in row]
            if not any(values):
                continue
            rows.append(values)
            max_cols = max(max_cols, len(values))
        sheets.append(
            {
                "title": ws.title,
                "rows": len(rows),
                "cols": max_cols,
                "preview": rows[:10],
            }
        )
    return {"sheets": sheets}


def extract_pptx_template_outline(filepath: str) -> dict[str, Any]:
    try:
        from pptx import Presentation
    except ImportError:
        logger.warning("python-pptx not installed; pptx outline extraction skipped")
        return {"slides": []}

    try:
        prs = Presentation(filepath)
    except Exception as exc:
        logger.warning("Failed to open PPTX template %s: %s", filepath, exc)
        return {"slides": []}

    slides: list[dict[str, Any]] = []
    for idx, slide in enumerate(prs.slides, start=1):
        texts: list[str] = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text = getattr(shape, "text", "").strip()
                if text:
                    texts.append(text)
        slides.append(
            {
                "index": idx,
                "texts": texts[:20],
                "shape_count": len(slide.shapes),
            }
        )
    return {"slides": slides}


def _build_docx(filepath: str, content: str):
    """Build a DOCX using OfficeCLI add command."""
    paragraphs = [p.strip() for p in content.split("\n\n") if p.strip()]
    for i, para in enumerate(paragraphs):
        if i == 0:
            _cli(["add", filepath, "/body", "--type", "paragraph",
                  "--prop", f"text={para}", "--prop", "style=Heading1"])
        elif len(para) < 60 and not para.endswith("。") and not para.endswith(")"):
            _cli(["add", filepath, "/body", "--type", "paragraph",
                  "--prop", f"text={para}", "--prop", "style=Heading2"])
        else:
            _cli(["add", filepath, "/body", "--type", "paragraph", "--prop", f"text={para}"])
    # Must close to persist changes to disk
    _cli(["close", filepath])


def _build_docx_from_template(template_path: str, output_path: str, content: str) -> None:
    """Copy a template docx and refill its body while keeping template styles/sections."""
    try:
        from docx import Document
    except ImportError:
        shutil.copyfile(template_path, output_path)
        _build_docx(output_path, content)
        return

    shutil.copyfile(template_path, output_path)
    doc = Document(output_path)
    template_paragraphs = [_capture_template_paragraph(p) for p in doc.paragraphs if (p.text or "").strip()]
    style_sequence = [p["style"] for p in template_paragraphs]
    heading_styles = [style for style in style_sequence if style.lower().startswith("heading")]
    body_style = next((style for style in style_sequence if not style.lower().startswith("heading")), "Normal")
    primary_heading = heading_styles[0] if heading_styles else "Heading 1"
    secondary_heading = heading_styles[1] if len(heading_styles) > 1 else primary_heading
    template_tables = [table for table in doc.tables]
    front_matter_elements = _capture_front_matter_elements(doc)
    preserved_heading_titles = {
        _normalize_heading(item.get("text", ""))
        for item in front_matter_elements
        if item["kind"] == "paragraph" and _should_preserve_template_table(item.get("text", ""))
    }
    preserved_table_indices: dict[str, int] = {}
    front_table_index = 0
    for item in front_matter_elements:
        if item["kind"] != "table":
            continue
        heading_text = _normalize_heading(item.get("heading", ""))
        if _should_preserve_template_table(heading_text):
            preserved_table_indices[heading_text] = front_table_index
        front_table_index += 1

    _clear_document_body(doc)
    _append_front_matter(doc, front_matter_elements)

    blocks = _parse_docx_blocks(content)
    paragraph_index = len([p for p in front_matter_elements if p["kind"] == "paragraph"])
    last_heading_text = ""
    for block in blocks:
        if block["type"] == "table":
            heading = _normalize_heading(last_heading_text)
            existing_table_index = preserved_table_indices.get(heading)
            if existing_table_index is not None and existing_table_index < len(doc.tables):
                _fill_table_cells(doc.tables[existing_table_index], block["rows"])
                continue
            _append_text_table(doc, block["rows"])
            continue

        para_text = block["text"]
        if block["type"] == "heading" and _normalize_heading(para_text) in preserved_heading_titles:
            last_heading_text = para_text
            continue
        para_template = _pick_template_paragraph(template_paragraphs, paragraph_index, para_text, primary_heading, secondary_heading, body_style)
        para = _append_paragraph_from_template(doc, para_template, para_text)
        if block["type"] == "heading":
            last_heading_text = para_text
        paragraph_index += 1

    doc.save(output_path)


def _pick_template_paragraph(
    template_paragraphs: list[dict[str, Any]],
    index: int,
    para_text: str,
    primary_heading: str,
    secondary_heading: str,
    body_style: str,
) -> dict[str, Any]:
    if index < len(template_paragraphs):
        return template_paragraphs[index]
    if index == 0:
        return {"style": primary_heading, "run": {}}
    if len(para_text) < 60 and not para_text.endswith("。") and not para_text.endswith(")"):
        return {"style": secondary_heading, "run": {}}
    return {"style": body_style, "run": {}}


def _clear_document_body(doc: Any) -> None:
    body = doc._element.body
    for child in list(body):
        if child.tag.endswith("sectPr"):
            continue
        body.remove(child)


def _insert_body_element(doc: Any, element: Any) -> None:
    body = doc._element.body
    children = list(body)
    sectpr_index = next((idx for idx, child in enumerate(children) if child.tag.endswith("sectPr")), None)
    if sectpr_index is None:
        body.append(element)
        return
    body.insert(sectpr_index, element)


def _capture_front_matter_elements(doc: Any) -> list[dict[str, Any]]:
    front_matter: list[dict[str, Any]] = []
    para_iter = iter(doc.paragraphs)
    table_iter = iter(doc.tables)
    last_heading_text = ""
    for child in list(doc._element.body):
        if child.tag.endswith("sectPr"):
            continue
        local = child.tag.split("}")[-1]
        if local == "p":
            para = next(para_iter, None)
            text = (para.text or "").strip() if para is not None else ""
            if _is_business_heading(text):
                break
            if text:
                last_heading_text = text
            front_matter.append({"kind": "paragraph", "element": deepcopy(child), "text": text})
        elif local == "tbl":
            table = next(table_iter, None)
            front_matter.append({"kind": "table", "element": deepcopy(child), "table": table, "heading": last_heading_text})
    return front_matter


def _append_front_matter(doc: Any, elements: list[dict[str, Any]]) -> None:
    for item in elements:
        _insert_body_element(doc, deepcopy(item["element"]))


def _parse_docx_blocks(content: str) -> list[dict[str, Any]]:
    lines = [line.rstrip() for line in content.splitlines()]
    blocks: list[dict[str, Any]] = []
    i = 0
    while i < len(lines):
        line = lines[i].strip()
        if not line:
            i += 1
            continue
        if line.startswith("|") and i + 1 < len(lines) and re.match(r"^\|?[-:\s|]+\|?$", lines[i + 1].strip()):
            table_lines = [line, lines[i + 1].strip()]
            i += 2
            while i < len(lines) and lines[i].strip().startswith("|"):
                table_lines.append(lines[i].strip())
                i += 1
            blocks.append({"type": "table", "rows": [_split_markdown_row(row) for row in [table_lines[0], *table_lines[2:]]]})
            continue
        if line.startswith("#") or _is_business_heading(line):
            blocks.append({"type": "heading", "text": _strip_heading_marker(line)})
        else:
            blocks.append({"type": "paragraph", "text": line})
        i += 1
    return blocks


def _looks_heading(text: str) -> bool:
    stripped = text.strip()
    return stripped.startswith("#") or bool(re.match(r"^(第[一二三四五六七八九十百零〇\d]+[章节条]|[一二三四五六七八九十]+[、\.])", stripped))


def _strip_heading_marker(text: str) -> str:
    stripped = text.strip()
    if stripped.startswith("#"):
        return stripped.lstrip("#").strip()
    return stripped


def _split_markdown_row(row: str) -> list[str]:
    parts = [part.strip() for part in row.strip().strip("|").split("|")]
    return parts


def _should_preserve_template_table(heading_text: str) -> bool:
    normalized = _normalize_heading(heading_text)
    preserve_titles = {"文档发行审批", "文件修订履历"}
    return normalized in preserve_titles


def _normalize_heading(text: str) -> str:
    return re.sub(r"^[#\s]+", "", (text or "").strip())


def _is_business_heading(text: str) -> bool:
    stripped = _normalize_heading(text)
    return bool(re.match(r"^(第[一二三四五六七八九十百零〇\d]+[章节条]|[一二三四五六七八九十]+[、\.]|[0-9]+[\.、])", stripped))


def _append_existing_table(doc: Any, template_table: Any, rows: list[list[str]]) -> None:
    tbl = deepcopy(template_table._element)
    doc._element.body.append(tbl)
    new_table = doc.tables[-1]
    _fill_table_cells(new_table, rows)


def _append_text_table(doc: Any, rows: list[list[str]]) -> None:
    if not rows:
        return
    table = doc.add_table(rows=len(rows), cols=max(len(r) for r in rows))
    _fill_table_cells(table, rows)


def _fill_table_cells(table: Any, rows: list[list[str]]) -> None:
    for r_idx, row in enumerate(rows):
        if r_idx >= len(table.rows):
            break
        for c_idx, cell_text in enumerate(row):
            if c_idx >= len(table.rows[r_idx].cells):
                break
            table.rows[r_idx].cells[c_idx].text = cell_text


def _capture_template_paragraph(paragraph: Any) -> dict[str, Any]:
    style_name = getattr(getattr(paragraph, "style", None), "name", "") or "Normal"
    run_info: dict[str, Any] = {}
    if paragraph.runs:
        run = paragraph.runs[0]
        font = getattr(run, "font", None)
        run_info = {
            "bold": run.bold,
            "italic": run.italic,
            "underline": run.underline,
            "font_name": getattr(font, "name", None) if font else None,
            "font_size": getattr(font, "size", None) if font else None,
        }
    return {"style": style_name, "run": run_info}


def _apply_run_formatting(paragraph: Any, template_info: dict[str, Any]) -> None:
    if not paragraph.runs:
        return
    run = paragraph.runs[0]
    info = template_info.get("run", {}) if isinstance(template_info, dict) else {}
    if not info:
        return
    run.bold = info.get("bold")
    run.italic = info.get("italic")
    run.underline = info.get("underline")
    font = getattr(run, "font", None)
    if font is not None:
        if info.get("font_name"):
            font.name = info["font_name"]
        if info.get("font_size") is not None:
            font.size = info["font_size"]


def _looks_heading(text: str) -> bool:
    stripped = text.strip()
    return stripped.startswith("#") or _is_business_heading(stripped)


def _should_preserve_template_table(heading_text: str) -> bool:
    normalized = _normalize_heading(heading_text)
    return normalized in {"文档发行审批", "文件修订履历"}


def _is_business_heading(text: str) -> bool:
    stripped = _normalize_heading(text)
    return bool(
        re.match(r"^第[一二三四五六七八九十百千万0-9]+[章节条款]", stripped)
        or re.match(r"^[一二三四五六七八九十]+[、.]", stripped)
        or re.match(r"^[0-9]+[、.]", stripped)
        or re.match(r"^[0-9]+\s*[\u4e00-\u9fffA-Za-z]", stripped)
    )


def _capture_template_paragraph(paragraph: Any) -> dict[str, Any]:
    style_name = getattr(getattr(paragraph, "style", None), "name", "") or "Normal"
    run_info: dict[str, Any] = {}
    if paragraph.runs:
        run = paragraph.runs[0]
        font = getattr(run, "font", None)
        run_info = {
            "bold": run.bold,
            "italic": run.italic,
            "underline": run.underline,
            "font_name": getattr(font, "name", None) if font else None,
            "font_size": getattr(font, "size", None) if font else None,
        }
    return {"style": style_name, "run": run_info, "element": deepcopy(paragraph._element)}


def _set_paragraph_text(paragraph: Any, text: str) -> None:
    if paragraph.runs:
        paragraph.runs[0].text = text
        for run in paragraph.runs[1:]:
            run.text = ""
        return
    paragraph.add_run(text)


def _append_paragraph_from_template(doc: Any, template_info: dict[str, Any], text: str) -> Any:
    from docx.text.paragraph import Paragraph

    template_element = template_info.get("element") if isinstance(template_info, dict) else None
    if template_element is None:
        para = doc.add_paragraph(text)
        try:
            para.style = template_info.get("style", "Normal")
        except Exception:
            pass
        _apply_run_formatting(para, template_info)
        return para
    new_element = deepcopy(template_element)
    _insert_body_element(doc, new_element)
    para = Paragraph(new_element, doc._body)
    _set_paragraph_text(para, text)
    _apply_run_formatting(para, template_info)
    return para


def _append_existing_table(doc: Any, template_table: Any, rows: list[list[str]]) -> None:
    tbl = deepcopy(template_table._element)
    _insert_body_element(doc, tbl)
    new_table = doc.tables[-1]
    _fill_table_cells(new_table, rows)


def _build_xlsx(filepath: str, content: str):
    """Build an XLSX using OfficeCLI."""
    rows = [line.split("\t") for line in content.strip().split("\n") if line.strip()]
    sheet = "/sheet1"
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell_ref = f"{sheet}/{chr(65 + ci)}{ri + 1}"
            _cli(["set", filepath, cell_ref, "--prop", f"value={val.strip()}"])
    _cli(["close", filepath])


def _build_xlsx_from_template(template_path: str, output_path: str, content: str) -> None:
    try:
        from openpyxl import load_workbook
    except ImportError:
        shutil.copyfile(template_path, output_path)
        _build_xlsx(output_path, content)
        return

    shutil.copyfile(template_path, output_path)
    wb = load_workbook(output_path)
    ws = wb.worksheets[0]
    rows = [line.split("\t") for line in content.strip().split("\n") if line.strip()]
    for ri, row in enumerate(rows, start=1):
        for ci, val in enumerate(row, start=1):
            ws.cell(row=ri, column=ci, value=val.strip())
    wb.save(output_path)


def _build_pptx(filepath: str, content: str):
    """Build a PPTX using OfficeCLI."""
    paragraphs = [p.strip() for p in content.split("\n\n") if p.strip()]
    for i, para in enumerate(paragraphs):
        if i == 0:
            _cli(["set", filepath, "/slide[1]", "--prop", f"title={para[:100]}"])
        else:
            _cli(["add", filepath, "/", "--type", "slide", "--prop", f"title={para[:100]}"])
    _cli(["close", filepath])


def _build_pptx_from_template(template_path: str, output_path: str, content: str) -> None:
    shutil.copyfile(template_path, output_path)
    _build_pptx(output_path, content)


def _sanitize_filename(name: str) -> str:
    return "".join(c if c.isalnum() or c in " _-." else "_" for c in name).strip() or "report"
