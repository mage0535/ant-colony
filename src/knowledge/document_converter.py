"""Multi-level document converter.

Extraction strategy levels:
  Level 1 — Remote URLs via firecrawl / jina / direct web fetch
  Level 2 — PyMuPDF (pymupdf) for text-based PDFs
  Level 3 — marker-pdf for scanned PDFs / OCR / complex layouts
  Level 4 — python-docx / python-pptx / openpyxl for Office formats
"""
from __future__ import annotations

import logging
import os
import urllib.parse
from pathlib import Path

logger = logging.getLogger(__name__)

_IMAGE_EXTENSIONS = {".jpg", ".jpeg", ".png", ".bmp", ".tiff", ".webp", ".gif"}


class DocumentConverter:
    """Multi-level document extraction with smart fallback between backends."""

    @staticmethod
    def _is_url(path_or_url: str) -> bool:
        return urllib.parse.urlparse(path_or_url).scheme in ("http", "https")

    @staticmethod
    def _guess_file_type(path: str) -> str:
        _, ext = os.path.splitext(path)
        return ext.lower()

    # ------------------------------------------------------------------
    # Public API
    # ------------------------------------------------------------------

    def convert(self, path_or_url: str, file_type: str | None = None) -> str:
        """Auto-detect source type and extract text content.

        Parameters
        ----------
        path_or_url:
            Local file path or remote URL.
        file_type:
            Optional override (e.g. ``"pdf"``, ``"docx"``).  Auto-detected
            from extension when omitted.

        Returns
        -------
        Extracted text (Markdown-ish for Office files, plain text for PDF).
        """
        if self._is_url(path_or_url):
            return self._fetch_url(path_or_url)

        filepath = Path(path_or_url)
        if not filepath.exists():
            raise FileNotFoundError(f"File not found: {path_or_url}")

        ext = (f".{file_type.lstrip('.')}" if file_type else filepath.suffix).lower()

        dispatch = {
            ".pdf": self.convert_pdf,
            ".docx": self.convert_docx,
            ".pptx": self.convert_pptx,
            ".xlsx": self.convert_xlsx,
        }
        handler = dispatch.get(ext)
        if handler is not None:
            return handler(str(filepath))
        if ext in _IMAGE_EXTENSIONS:
            return self.convert_image(str(filepath))
        return filepath.read_text(encoding="utf-8", errors="replace")

    def convert_pdf(self, path: str, ocr: bool = False) -> str:
        """Extract text from a PDF.

        Strategy
        --------
        1.  **PyMuPDF** (fast, text-based PDFs).
        2.  If result is empty or suspiciously short for the file size and
            *ocr=True* is passed, fall back to **marker-pdf**.
        3.  If *ocr=False* but the result looks like a scanned document,
            a warning is logged; the caller can re-try with ``ocr=True``.
        """
        file_size = os.path.getsize(path)

        # -- Level 2: PyMuPDF ------------------------------------------------
        text = self._extract_pymupdf(path)
        if text and self._is_plausible(text, file_size):
            return text

        # -- Level 3: marker-pdf fallback ------------------------------------
        if not self._is_plausible(text, file_size):
            if ocr:
                return self._extract_marker_pdf(path)
            logger.warning(
                "PyMuPDF output for %s seems incomplete (%d chars for %d byte file). "
                "Re-run with ocr=True to enable marker-pdf OCR.",
                path, len(text or ""), file_size,
            )
        return text or ""

    def convert_docx(self, path: str) -> str:
        """Extract text from a DOCX file via python-docx."""
        return self._extract_docx(path)

    def convert_pptx(self, path: str) -> str:
        """Extract text from a PPTX file via python-pptx."""
        return self._extract_pptx(path)

    def convert_xlsx(self, path: str) -> str:
        """Extract text from an XLSX file via openpyxl."""
        return self._extract_xlsx(path)

    def convert_image(self, path: str) -> str:
        """OCR an image via marker-pdf (if available)."""
        return self._extract_image_ocr(path)

    # ------------------------------------------------------------------
    # Level 1 — Remote URL fetching
    # ------------------------------------------------------------------

    def _fetch_url(self, url: str) -> str:
        """Attempt firecrawl → jina reader → direct HTTP GET."""
        # -- firecrawl -------------------------------------------------------
        try:
            return self._fetch_via_firecrawl(url)
        except Exception as exc:
            logger.debug("firecrawl failed for %s: %s", url, exc)

        # -- jina reader -----------------------------------------------------
        try:
            return self._fetch_via_jina(url)
        except Exception as exc:
            logger.debug("jina reader failed for %s: %s", url, exc)

        # -- direct fetch ----------------------------------------------------
        return self._fetch_direct(url)

    @staticmethod
    def _fetch_via_firecrawl(url: str) -> str:
        """Extract via firecrawl Python SDK (firecrawl-py)."""
        try:
            from firecrawl import FirecrawlApp
        except ImportError:
            raise ImportError("firecrawl-py not installed")

        app = FirecrawlApp()
        result = app.scrape_url(url)
        return result.get("content") or result.get("markdown") or ""

    @staticmethod
    def _fetch_via_jina(url: str) -> str:
        """Extract via jina.ai reader API."""
        import requests
        jina_url = f"https://r.jina.ai/{url}"
        resp = requests.get(jina_url, timeout=30)
        resp.raise_for_status()
        text = resp.text
        if text.startswith("data:"):
            lines = [l for l in text.splitlines() if not l.startswith("data:")]
            text = "\n".join(lines)
        return text

    @staticmethod
    def _fetch_direct(url: str) -> str:
        """Fallback direct HTML fetch with best-effort text extraction."""
        import requests
        from bs4 import BeautifulSoup

        resp = requests.get(url, timeout=30, headers={
            "User-Agent": (
                "Mozilla/5.0 (Windows NT 10.0; Win64; x64) "
                "AppleWebKit/537.36 (KHTML, like Gecko) "
                "Chrome/120.0.0.0 Safari/537.36"
            ),
        })
        resp.raise_for_status()
        ct = resp.headers.get("content-type", "")
        if "text/plain" in ct:
            return resp.text
        soup = BeautifulSoup(resp.text, "html.parser")
        for tag in soup(["script", "style", "nav", "footer", "header"]):
            tag.decompose()
        return soup.get_text(separator="\n", strip=True)

    # ------------------------------------------------------------------
    # Level 2 — PyMuPDF
    # ------------------------------------------------------------------

    @staticmethod
    def _extract_pymupdf(path: str) -> str:
        """Extract text from a PDF using PyMuPDF."""
        try:
            import pymupdf
        except ImportError:
            logger.warning("PyMuPDF (pymupdf) not installed — cannot extract PDF")
            return ""

        pages: list[str] = []
        try:
            doc = pymupdf.open(path)
        except Exception as exc:
            logger.warning("PyMuPDF failed to open %s: %s", path, exc)
            return ""

        for page_num in range(doc.page_count):
            try:
                page = doc.load_page(page_num)
                text = page.get_text()
                if text:
                    pages.append(text)
            except Exception as exc:
                logger.debug("PyMuPDF page %d error: %s", page_num, exc)
        doc.close()
        return "\n\n".join(pages)

    # ------------------------------------------------------------------
    # Level 3 — marker-pdf (OCR / complex layout)
    # ------------------------------------------------------------------

    @staticmethod
    def _extract_marker_pdf(path: str) -> str:
        """Extract text from a PDF using marker-pdf (OCR pipeline)."""
        try:
            from marker.converters.pdf import PdfConverter
            from marker.models import create_model_dict
            from marker.output import text_from_output
        except ImportError:
            logger.warning(
                "marker-pdf not installed — cannot OCR %s. "
                "Install with: pip install marker-pdf",
                path,
            )
            return ""

        try:
            converter = PdfConverter(
                artifact_dict=create_model_dict(),
            )
            rendered = converter(path)
            return text_from_output(rendered)
        except Exception as exc:
            logger.warning("marker-pdf conversion failed for %s: %s", path, exc)
            return ""

    # ------------------------------------------------------------------
    # Level 4 — Office formats
    # ------------------------------------------------------------------

    @staticmethod
    def _extract_docx(path: str) -> str:
        """Extract text from a DOCX file."""
        try:
            from docx import Document
        except ImportError:
            logger.warning("python-docx not installed")
            return ""

        try:
            doc = Document(path)
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            return "\n\n".join(paragraphs)
        except Exception as exc:
            logger.warning("Failed to read DOCX %s: %s", path, exc)
            return ""

    @staticmethod
    def _extract_pptx(path: str) -> str:
        """Extract text from a PPTX file."""
        try:
            from pptx import Presentation
        except ImportError:
            logger.warning("python-pptx not installed")
            return ""

        parts: list[str] = []
        try:
            prs = Presentation(path)
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_texts: list[str] = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            t = paragraph.text.strip()
                            if t:
                                slide_texts.append(t)
                    if shape.has_table:
                        for row in shape.table.rows:
                            cells = [cell.text.strip() for cell in row.cells]
                            slide_texts.append(" | ".join(cells))
                if slide_texts:
                    parts.append(f"--- Slide {slide_num} ---\n" + "\n".join(slide_texts))
        except Exception as exc:
            logger.warning("Failed to read PPTX %s: %s", path, exc)
            return ""
        return "\n\n".join(parts)

    @staticmethod
    def _extract_xlsx(path: str) -> str:
        """Extract text from an XLSX file."""
        try:
            import openpyxl
        except ImportError:
            logger.warning("openpyxl not installed")
            return ""

        parts: list[str] = []
        try:
            wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                rows: list[str] = []
                for row in ws.iter_rows(values_only=True):
                    cells = [str(c) if c is not None else "" for c in row]
                    line = "\t".join(cells)
                    if line.strip():
                        rows.append(line)
                if rows:
                    parts.append(f"=== {sheet_name} ===\n" + "\n".join(rows))
            wb.close()
        except Exception as exc:
            logger.warning("Failed to read XLSX %s: %s", path, exc)
            return ""
        return "\n\n".join(parts)

    # ------------------------------------------------------------------
    # Image OCR helper
    # ------------------------------------------------------------------

    @staticmethod
    def _extract_image_ocr(path: str) -> str:
        """OCR a single image via marker-pdf (if installed)."""
        try:
            from marker.converters.image import ImageConverter
            from marker.models import create_model_dict
        except ImportError:
            logger.warning(
                "marker-pdf not installed — cannot OCR image %s", path
            )
            return ""

        try:
            converter = ImageConverter(artifact_dict=create_model_dict())
            rendered = converter(path)
            return rendered.text if hasattr(rendered, "text") else str(rendered)
        except Exception as exc:
            logger.warning("Image OCR failed for %s: %s", path, exc)
            return ""

    # ------------------------------------------------------------------
    # Heuristics
    # ------------------------------------------------------------------

    @staticmethod
    def _is_plausible(text: str | None, file_size: int) -> bool:
        """Heuristic: does extracted text plausibly represent the file?"""
        if not text or not text.strip():
            return False
        text_len = len(text)
        return text_len > max(50, file_size * 0.02)


# Module-level wrappers for backward compatibility with collector imports
_converter = DocumentConverter()


def convert_document(path: str) -> str | None:
    """Convert a document file to text. Wraps DocumentConverter.convert()."""
    return _converter.convert(path)


def guess_type(path: str) -> str:
    """Guess the file type from extension. Wraps DocumentConverter._guess_file_type()."""
    return _converter._guess_file_type(path)
